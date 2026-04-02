# -*- coding: utf-8 -*-
"""
CSV（Headline, Subheadline, Body, CTA, Hook, VisualIdea）から PowerPoint を生成する。
Canva へは「ファイル → インポート」または pptx をアップロードして編集する想定。

使い方:
  python csv_to_pptx.py -i canva_part16_最幸ママ_bulk_only.csv -o output.pptx
  python csv_to_pptx.py -i canva_part19_最幸ママ_bulk_only.csv -o パート19.pptx --title "最幸ママ・アカデミー"
"""
from __future__ import annotations

import argparse
import csv
import sys
from io import StringIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN


EXPECTED_COLUMNS = ("Headline", "Subheadline", "Body", "CTA", "Hook", "VisualIdea")


def read_rows(csv_path: Path) -> list[dict[str, str]]:
    text = csv_path.read_text(encoding="utf-8-sig")
    if not text.strip():
        return []
    reader = csv.DictReader(StringIO(text))
    rows = []
    for row in reader:
        if row is None:
            continue
        cleaned = {k: (v or "").strip() for k, v in row.items()}
        if not any(cleaned.values()):
            continue
        rows.append(cleaned)
    return rows


def build_body_text(row: dict[str, str]) -> str:
    parts: list[str] = []
    if row.get("Subheadline"):
        parts.append(row["Subheadline"])
    if row.get("Body"):
        parts.append(row["Body"])
    extras = []
    if row.get("CTA"):
        extras.append(f"CTA：{row['CTA']}")
    if row.get("Hook"):
        extras.append(f"Hook：{row['Hook']}")
    if row.get("VisualIdea"):
        extras.append(f"ビジュアル：{row['VisualIdea']}")
    if extras:
        parts.append("\n".join(extras))
    return "\n\n".join(parts) if parts else ""


def set_title_font(title_shape, size_pt: float = 28):
    if not title_shape.has_text_frame:
        return
    for p in title_shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size_pt)


def set_body_font(body_shape, base_pt: float = 18):
    if not body_shape.has_text_frame:
        return
    tf = body_shape.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(base_pt)


def add_content_slide(prs: Presentation, headline: str, body: str):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    body_ph = slide.placeholders[1]

    title.text = headline or "（タイトルなし）"
    set_title_font(title, 26)

    tf = body_ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = body or ""
    p.alignment = PP_ALIGN.LEFT
    set_body_font(body_ph, 16)

    # タイトルが長い場合は少し小さく
    if headline and len(headline) > 40:
        set_title_font(title, 22)


def run(input_csv: Path, output_pptx: Path, deck_title: str | None) -> None:
    rows = read_rows(input_csv)
    if not rows:
        print("No data rows found in CSV.", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()

    # 表紙（任意）
    if deck_title:
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = deck_title
        if slide.placeholders and len(slide.placeholders) > 1:
            try:
                slide.placeholders[1].text_frame.text = input_csv.stem
            except Exception:
                pass

    for row in rows:
        headline = row.get("Headline", "")
        body = build_body_text(row)
        add_content_slide(prs, headline, body)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))
    print(f"Saved: {output_pptx} ({len(rows)} content slides)")


def main() -> None:
    parser = argparse.ArgumentParser(description="Build PPTX from Canva-style CSV.")
    parser.add_argument("-i", "--input", required=True, type=Path, help="Input CSV path")
    parser.add_argument("-o", "--output", required=True, type=Path, help="Output .pptx path")
    parser.add_argument(
        "--title",
        default="",
        help="Optional cover slide title (deck title). Empty = no cover slide.",
    )
    args = parser.parse_args()

    if not args.input.is_file():
        print(f"Input not found: {args.input}", file=sys.stderr)
        sys.exit(1)

    run(args.input, args.output, args.title.strip() or None)


if __name__ == "__main__":
    main()
